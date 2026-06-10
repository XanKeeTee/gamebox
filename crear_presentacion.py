from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Inicializar presentación
prs = Presentation()

# Configurar relación de aspecto a 16:9
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Paleta de colores Gamebox
bg_color = RGBColor(10, 10, 12)       # Fondo oscuro
text_white = RGBColor(255, 255, 255)  
text_cyan = RGBColor(6, 182, 212)     # Cian neón
text_fuchsia = RGBColor(217, 70, 239) # Fucsia neón
text_gray = RGBColor(156, 163, 175)   

def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    return shape

def add_title(slide, text, color=text_white):
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial Black"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.italic = True
    p.font.color.rgb = color

def add_bullet(tf, text, size=32, color=text_gray, bold=False):
    p = tf.add_paragraph()
    p.text = "• " + text
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(20)

# ==========================================
# DIAPOSITIVA 1: Portada
# ==========================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1)
txBox = slide1.shapes.add_textbox(Inches(0), Inches(3), prs.slide_width, Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "GAMEBOX"
p.font.name = "Arial Black"
p.font.size = Pt(110)
p.font.italic = True
p.font.color.rgb = text_white
p.alignment = PP_ALIGN.CENTER

txBox2 = slide1.shapes.add_textbox(Inches(0), Inches(5.2), prs.slide_width, Inches(1))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "TU RED SOCIAL & GESTOR DE VIDEOJUEGOS"
p2.font.name = "Arial"
p2.font.size = Pt(28)
p2.font.bold = True
p2.font.color.rgb = text_cyan
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide1.shapes.add_textbox(Inches(0), Inches(7.5), prs.slide_width, Inches(1))
p3 = txBox3.text_frame.paragraphs[0]
p3.text = "Desarrollado por Miguel Ángel Gambero | 2º DAW"
p3.font.name = "Arial"
p3.font.size = Pt(18)
p3.font.color.rgb = text_gray
p3.alignment = PP_ALIGN.CENTER

# ==========================================
# DIAPOSITIVA 2: El Problema
# ==========================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_title(slide2, "EL PROBLEMA", text_fuchsia)
txBox = slide2.shapes.add_textbox(Inches(1.5), Inches(3), Inches(13), Inches(5))
tf = txBox.text_frame
add_bullet(tf, "Información gamer fragmentada (Steam, PS, Xbox, Nintendo).", 36, text_white)
add_bullet(tf, "No existe un lugar unificado y visual para gestionar el 'Backlog'.", 36, text_white)
add_bullet(tf, "Ausencia de un ecosistema verdaderamente social y moderno.", 36, text_white)

# ==========================================
# DIAPOSITIVA 3: La Solución
# ==========================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_title(slide3, "LA SOLUCIÓN: EL HUB DEFINITIVO", text_cyan)
txBox = slide3.shapes.add_textbox(Inches(1.5), Inches(3), Inches(13), Inches(5))
tf = txBox.text_frame
add_bullet(tf, "Gestión total de la colección (Jugando, Terminado, Abandonado).", 32, text_gray)
add_bullet(tf, "Integración y sincronización en tiempo real con la API de IGDB.", 32, text_gray)
add_bullet(tf, "Muro de actividad (Feed) para compartir reseñas y capturas.", 32, text_gray)

# ==========================================
# DIAPOSITIVA 4: Arquitectura
# ==========================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_title(slide4, "ARQUITECTURA DEL PROYECTO", text_white)

s1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3.5), Inches(4), Inches(4))
s1.fill.solid()
s1.fill.fore_color.rgb = RGBColor(22, 27, 34)
s1.line.color.rgb = text_cyan
s1.line.width = Pt(3)
p1 = s1.text_frame.paragraphs[0]
p1.text = "API IGDB\n\nConsultas asíncronas en tiempo real. Catálogo infinito sin saturar la BD."
p1.font.size = Pt(24)
p1.font.color.rgb = text_white

s2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(3.5), Inches(4), Inches(4))
s2.fill.solid()
s2.fill.fore_color.rgb = RGBColor(22, 27, 34)
s2.line.color.rgb = text_fuchsia
s2.line.width = Pt(3)
p2 = s2.text_frame.paragraphs[0]
p2.text = "CACHÉ AVANZADA\n\nAlmacenamiento temporal para tiempos de carga ultrarrápidos."
p2.font.size = Pt(24)
p2.font.color.rgb = text_white

s3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11), Inches(3.5), Inches(4), Inches(4))
s3.fill.solid()
s3.fill.fore_color.rgb = RGBColor(22, 27, 34)
s3.line.color.rgb = text_white
s3.line.width = Pt(3)
p3 = s3.text_frame.paragraphs[0]
p3.text = "DJANGO ORM\n\nBase de datos robusta (PostgreSQL) para lógica social."
p3.font.size = Pt(24)
p3.font.color.rgb = text_white

# ==========================================
# DIAPOSITIVA 5: Stack Tecnológico
# ==========================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_title(slide5, "STACK TECNOLÓGICO", text_fuchsia)
txBox = slide5.shapes.add_textbox(Inches(1.5), Inches(3), Inches(13), Inches(5))
tf = txBox.text_frame
add_bullet(tf, "Backend: Python 3 y Django Framework", 36, text_white, True)
add_bullet(tf, "Base de Datos: PostgreSQL (Alojado en Neon.tech)", 36, text_white, True)
add_bullet(tf, "Frontend: HTML5, Tailwind CSS y JavaScript Puro", 36, text_white, True)
add_bullet(tf, "Multimedia: Cloudinary (SaaS de almacenamiento)", 36, text_white, True)

# ==========================================
# DIAPOSITIVA 6: Interfaz Visual (Placeholders)
# ==========================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)
add_title(slide6, "INTERFAZ Y UX PREMIUM", text_cyan)

b1 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(4.2), Inches(4.5))
b1.fill.solid()
b1.fill.fore_color.rgb = RGBColor(22, 27, 34)
p = b1.text_frame.paragraphs[0]
p.text = "[Añade aquí captura del FEED]"
p.font.color.rgb = text_gray

b2 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.9), Inches(2.5), Inches(4.2), Inches(5.5))
b2.fill.solid()
b2.fill.fore_color.rgb = RGBColor(22, 27, 34)
b2.line.color.rgb = text_fuchsia
b2.line.width = Pt(3)
p = b2.text_frame.paragraphs[0]
p.text = "[Añade aquí captura del PERFIL]"
p.font.color.rgb = text_fuchsia

b3 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.8), Inches(3), Inches(4.2), Inches(4.5))
b3.fill.solid()
b3.fill.fore_color.rgb = RGBColor(22, 27, 34)
p = b3.text_frame.paragraphs[0]
p.text = "[Añade aquí captura de LA COLECCIÓN]"
p.font.color.rgb = text_gray

# ==========================================
# DIAPOSITIVA 7: Cierre
# ==========================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)
txBox = slide7.shapes.add_textbox(Inches(0), Inches(3), prs.slide_width, Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "EL PROYECTO ESTÁ VIVO"
p.font.name = "Arial Black"
p.font.size = Pt(72)
p.font.italic = True
p.font.color.rgb = text_white
p.alignment = PP_ALIGN.CENTER

txBox2 = slide7.shapes.add_textbox(Inches(0), Inches(5.5), prs.slide_width, Inches(1))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Desplegado en producción con Render"
p2.font.name = "Arial"
p2.font.size = Pt(28)
p2.font.color.rgb = text_cyan
p2.alignment = PP_ALIGN.CENTER

prs.save('Presentacion_Gamebox.pptx')
print("¡Presentación creada con éxito! Abre Presentacion_Gamebox.pptx")